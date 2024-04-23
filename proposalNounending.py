import streamlit as st
from openai import OpenAI
import io
from PyPDF2 import PdfReader
from pptx import Presentation

st.title('Proposal Revision: Noun-ending')

def read_pdf(file):
    with io.BytesIO(file.getvalue()) as f:
        reader = PdfReader(f)
        text = ''
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text

def read_pptx(file):
    with io.BytesIO(file.getvalue()) as f:
        prs = Presentation(f)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

def make_guide_from_proposal(proposal_text):
    client = OpenAI()
    
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": """
            명사형 종결 처리가 안된 부분을 찾아 작성한 뒤,
            (추천)에 명사형 종결 처리된 문장으로 변경

            Format Example : 
            [1] 입찰을 제안드립니다. -> (추천) 입찰 제안
            
            [2] 본 프로젝트는 성남시의 디지털 트랜스포메이션을 목표로 합니다. -> (추천) 본 프로젝트는 성남시의 디지털 트랜스포메이션이 목표
            
            [3] 명시적으로 작성해 주세요. -> (추천) 명시적으로 작성
            
             """},
            {"role": "user", "content": proposal_text}
        ]
    )
    
    return response.choices[0].message.content

uploaded_file = st.file_uploader("작성하신 제안서를 업로드 하세요.", type=['pptx', 'pdf'])

if uploaded_file:
    if uploaded_file.type == "application/pdf":
        file_content = read_pdf(uploaded_file)
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        file_content = read_pptx(uploaded_file)
    else:
        st.error("지원하지 않는 파일 형식입니다. PDF 또는 PPTX 파일만 업로드 가능합니다.")
        file_content = None

    if file_content and st.button("제안서 가이드 받기"):
        try:
            guide = make_guide_from_proposal(file_content)
            st.write(guide)
        except Exception as e:
            st.error("Error: " + str(e))
