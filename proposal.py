import streamlit as st
import pandas as pd
import numpy as np
from openai import OpenAI
import io
from PyPDF2 import PdfReader
from pptx import Presentation

st.title('Proposal Revision: Terminology')

def read_pdf(file):
    with io.BytesIO(file.getvalue()) as f:
        reader = PdfReader(f)
        text = ''
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text

def read_pptx(file):
    with io.BytesIO(file.getvalue()) as f:
        prs = Presentation(f)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

def make_guide_from_proposal(proposal_text):
    client = OpenAI()
    
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": """
             이 내용에서 통일해야 하는 용어 후보를 리스트업 후 각각에 대해서 통일할 용어를 추천해줘. 추천 단어를 AI로 한정 짓는 건 아니고 일관성 면을 고려해서 추천해줘 한글로 작성 
            Answer Format Example ) 
            [1] 추천 단어: AI 
            - Term 1: 인공지능
            - Term 2: AI
            - 추천 이유: XXXXX
             """},
            {"role": "user", "content": proposal_text}
        ]
    )
    
    return response.choices[0].message.content

uploaded_file = st.file_uploader("작성하신 제안서를 업로드 하세요.", type=['pptx', 'pdf'])

if uploaded_file:
    if uploaded_file.type == "application/pdf":
        file_content = read_pdf(uploaded_file)
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        file_content = read_pptx(uploaded_file)
    else:
        st.error("지원하지 않는 파일 형식입니다. PDF 또는 PPTX 파일만 업로드 가능합니다.")
        file_content = None

    if file_content and st.button("제안서 가이드 받기"):
        try:
            guide = make_guide_from_proposal(file_content)
            st.write(guide)
        except Exception as e:
            st.error("Error: " + str(e))
